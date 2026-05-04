from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(34, 139, 34)  # Forest Green
SECONDARY_COLOR = RGBColor(70, 130, 180)  # Steel Blue
ACCENT_COLOR = RGBColor(255, 140, 0)  # Dark Orange
TEXT_COLOR = RGBColor(33, 33, 33)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_list):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Add title text
    title_frame = title_shape.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    title_frame = title_shape.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    # Left title
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    p.space_after = Pt(10)
    
    # Left items
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    # Right title
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    p.space_after = Pt(10)
    
    # Right items
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(6)
    
    return slide

# ============= SLIDE 1: Title Slide =============
add_title_slide(prs, "SVASTHYA FRESH", 
                "E-Commerce & Delivery Management System")

# ============= SLIDE 2: Project Overview =============
add_content_slide(prs, "Project Overview", [
    "Project Name: Svasthya Fresh",
    "Type: Full-Stack E-Commerce & Delivery Management Platform",
    "Purpose: Provide seamless ordering, delivery tracking, and admin management",
    "Target Users: End customers, Admin panel users, Support team",
    "Status: Under Development"
])

# ============= SLIDE 3: Project Goals =============
add_content_slide(prs, "Project Goals", [
    "Enable customers to browse and purchase products online",
    "Manage orders and track delivery in real-time",
    "Provide comprehensive admin dashboard for business analytics",
    "Support multiple payment methods and coupon management",
    "Implement robust authentication and security measures",
    "Deliver excellent customer support through ticketing system"
])

# ============= SLIDE 4: System Architecture =============
add_two_column_slide(prs, "System Architecture",
    "Frontend", [
        "Admin Dashboard (React + TypeScript)",
        "Customer Website (React + Vite)",
        "Component Library: shadcn-ui",
        "Styling: Tailwind CSS"
    ],
    "Backend", [
        "Spring Boot 3.5.10 (Java 17)",
        "REST API Architecture",
        "Database: MySQL",
        "Security: Spring Security"
    ])

# ============= SLIDE 5: Modules Overview =============
add_content_slide(prs, "Main Modules", [
    "1. User & Authentication Management",
    "2. Product & Inventory Management",
    "3. Order & Cart Management",
    "4. Payment & Transaction Processing",
    "5. Coupon & Discount Management",
    "6. Analytics & Reporting Dashboard",
    "7. Support Ticket System",
    "8. Wishlist & Personalization"
])

# ============= SLIDE 6: Module 1 - User Management =============
add_content_slide(prs, "Module 1: User & Authentication", [
    "Email & Mobile OTP Verification",
    "Google OAuth Authentication",
    "Admin User Registration & Login",
    "User Profile Management",
    "Password Reset Functionality",
    "Role-based Access Control (RBAC)",
    "Session Management with JWT Tokens",
    "User Query & Support Request Tracking"
])

# ============= SLIDE 7: Module 2 - Products & Categories =============
add_content_slide(prs, "Module 2: Products & Categories", [
    "Product Catalog Management",
    "Category Organization & Hierarchy",
    "Product Variants & Options",
    "Product Images & Media Gallery",
    "Price History Tracking",
    "Inventory Management",
    "Product Performance Analytics",
    "Banner Management for Promotions"
])

# ============= SLIDE 8: Module 3 - Orders & Cart =============
add_content_slide(prs, "Module 3: Orders & Cart Management", [
    "Shopping Cart Management",
    "Add/Remove/Update Cart Items",
    "Checkout Process",
    "Order Creation & Tracking",
    "Order Item Management",
    "Order Status Updates",
    "Order Details Viewing",
    "Order History Management"
])

# ============= SLIDE 9: Module 4 - Payments & Transactions =============
add_content_slide(prs, "Module 4: Payments & Transactions", [
    "Payment Processing Integration",
    "Multiple Payment Methods",
    "Transaction History Tracking",
    "Payment Refund Management",
    "Payment Reports & Analytics",
    "Transaction Verification",
    "Revenue Tracking",
    "Financial Dashboard"
])

# ============= SLIDE 10: Module 5 - Coupons & Discounts =============
add_content_slide(prs, "Module 5: Coupons & Discounts", [
    "Coupon Creation & Management",
    "Discount Code Generation",
    "Coupon Validation & Application",
    "Usage Tracking per Coupon",
    "Expiration Date Management",
    "Coupon Performance Reports",
    "Bulk Coupon Operations",
    "Customer Coupon History"
])

# ============= SLIDE 11: Module 6 - Analytics & Reports =============
add_content_slide(prs, "Module 6: Analytics & Reporting", [
    "Dashboard Analytics Overview",
    "Revenue Reports",
    "Sales Funnel Analysis",
    "Product Performance Metrics",
    "Customer Demographics Report",
    "Inventory Status Reports",
    "Order Status Analytics",
    "Payment & Refund Reports"
])

# ============= SLIDE 12: Module 7 - Support System =============
add_content_slide(prs, "Module 7: Support Ticket System", [
    "Create Support Tickets",
    "Ticket Status Management",
    "Message Threading",
    "Support Image Attachments",
    "Admin Support Dashboard",
    "Ticket Priority Levels",
    "Response Tracking",
    "Customer Support History"
])

# ============= SLIDE 13: Module 8 - Wishlist =============
add_content_slide(prs, "Module 8: Wishlist & Favorites", [
    "Add Products to Wishlist",
    "Wishlist Management",
    "Wishlist Sharing",
    "Quick Add to Cart from Wishlist",
    "Wishlist Item Tracking",
    "Price Drop Notifications",
    "Personalized Recommendations",
    "Wishlist Analytics"
])

# ============= SLIDE 14: Additional Features =============
add_content_slide(prs, "Additional Features", [
    "Address Management - Store multiple addresses",
    "Email & SMS Notifications",
    "OTP Rate Limiting",
    "Google OAuth Integration",
    "AWS S3 Integration for media storage",
    "SNS for push notifications",
    "Error Handling & Exception Management",
    "Health Check & API Monitoring"
])

# ============= SLIDE 15: Technology Stack =============
add_two_column_slide(prs, "Technology Stack - Frontend & Backend",
    "Frontend Technologies", [
        "React 18.2.0",
        "TypeScript",
        "Vite (Build Tool)",
        "Tailwind CSS",
        "shadcn-ui Components",
        "Radix UI",
        "React Hook Form",
        "TanStack React Query"
    ],
    "Backend Technologies", [
        "Spring Boot 3.5.10",
        "Java 17",
        "Spring Data JPA",
        "Spring Security",
        "MySQL Database",
        "JWT Authentication",
        "OpenAPI/Swagger",
        "Maven (Build Tool)"
    ])

# ============= SLIDE 16: Database Design - ER Diagram Tables =============
add_content_slide(prs, "Database Tables (ER Diagram)", [
    "Core Tables: Users, Roles, Permissions",
    "Product Management: Products, Categories, ProductVariants, ProductImages",
    "Order Processing: Orders, OrderItems, OrderEntity",
    "Shopping: Cart, CartItems, Wishlist, WishlistItems",
    "Financial: Payments, Coupons, CouponUsage, PriceHistory",
    "Authentication: Token, ActiveToken, EmailOtp, MobileOtp, PasswordReset",
    "Support: Support, SupportMessages, SupportImages",
    "Additional: Banners, Visitors, UserQueries, Addresses"
])

# ============= SLIDE 17: Data Flow Diagram (DFD) - Level 0 =============
add_content_slide(prs, "Data Flow Diagram (DFD) - Level 0", [
    "External Entities:",
    "  • Customer Users",
    "  • Admin Users",
    "  • Payment Gateway",
    "  • Email/SMS Service",
    "Main Process: E-Commerce & Delivery Management System",
    "Data Stores: Database",
    "External Systems: Google OAuth, AWS S3, AWS SNS"
])

# ============= SLIDE 18: Data Flow Diagram (DFD) - Level 1 =============
add_content_slide(prs, "Data Flow Diagram (DFD) - Level 1", [
    "1.0 Authentication & Authorization - User Login/Registration",
    "2.0 Product Browse & Search - Category/Product Catalog",
    "3.0 Cart & Order Management - Add items, Checkout",
    "4.0 Payment Processing - Transaction handling",
    "5.0 Order Fulfillment - Delivery tracking",
    "6.0 Analytics & Reporting - Business insights",
    "7.0 Support Management - Ticket handling"
])

# ============= SLIDE 19: Database Schema - Key Relationships =============
add_content_slide(prs, "Key Database Relationships", [
    "Users → Addresses (1:M) - Multiple addresses per user",
    "Users → Orders (1:M) - Customer orders",
    "Products → ProductVariants (1:M) - Product options",
    "Products → ProductImages (1:M) - Multiple images",
    "Orders → OrderItems (1:M) - Items in order",
    "Coupons → CouponUsage (1:M) - Usage tracking",
    "Users → WishlistItems (1:M) - Wishlist items",
    "Support → SupportMessages (1:M) - Ticket conversations"
])

# ============= SLIDE 20: Admin Dashboard Features =============
add_content_slide(prs, "Admin Dashboard Pages", [
    "Dashboard - Overview & KPIs",
    "Products - Manage catalog",
    "Categories - Organize products",
    "Coupons - Create discounts",
    "Orders - View & manage orders",
    "Payments - Transaction tracking",
    "Banners - Promotional content",
    "Members - User management",
    "Analytics - Reports & insights",
    "Support Center - Ticket management",
    "Settings - Configuration"
])

# ============= SLIDE 21: Task Assignment =============
add_content_slide(prs, "Task Assignments", [
    "Backend Development: API endpoints, Database operations",
    "Frontend Admin: Dashboard UI, Components, Forms",
    "Frontend Website: Customer portal, Shopping experience",
    "Authentication: Security implementation, OTP, JWT",
    "Payment Integration: Gateway setup, Transaction handling",
    "Analytics: Report generation, Data visualization",
    "Testing: Unit tests, Integration tests, E2E tests",
    "Deployment: Server setup, CI/CD pipeline, Monitoring"
])

# ============= SLIDE 22: Security Features =============
add_content_slide(prs, "Security Implementation", [
    "JWT Token-based Authentication",
    "Spring Security Framework",
    "Role-based Access Control (RBAC)",
    "Password Encryption",
    "OTP Verification (Email & SMS)",
    "Rate Limiting on OTP",
    "CORS Configuration",
    "Input Validation & Sanitization",
    "SQL Injection Prevention",
    "XSS Protection"
])

# ============= SLIDE 23: API Documentation =============
add_content_slide(prs, "API Documentation", [
    "OpenAPI/Swagger Integration",
    "RESTful API Endpoints",
    "Request/Response Models",
    "Authentication Headers",
    "Error Response Codes",
    "Pagination Support",
    "Filtering & Sorting",
    "API Rate Limiting"
])

# ============= SLIDE 24: Development Workflow =============
add_two_column_slide(prs, "Development & Deployment Workflow",
    "Development", [
        "Git version control",
        "Feature branches",
        "Code reviews",
        "Local testing",
        "Linting & formatting",
        "Build automation"
    ],
    "Deployment", [
        "Maven builds",
        "Docker containers",
        "Environment config",
        "Database migrations",
        "Monitoring & logs",
        "CI/CD pipeline"
    ])

# ============= SLIDE 25: Key Achievements =============
add_content_slide(prs, "Key Achievements", [
    "Complete API architecture implemented",
    "Responsive admin dashboard created",
    "Customer-facing website launched",
    "Multi-layer authentication system",
    "Comprehensive analytics engine",
    "Support ticket management system",
    "Payment processing integration",
    "Scalable database design"
])

# ============= SLIDE 26: Challenges & Solutions =============
add_two_column_slide(prs, "Challenges & Solutions",
    "Technical Challenges", [
        "Real-time order tracking",
        "Performance optimization",
        "Database scaling",
        "API rate limiting",
        "Payment gateway integration"
    ],
    "Solutions Implemented", [
        "WebSocket implementation",
        "Query optimization, caching",
        "Database indexing, partitioning",
        "Token bucket algorithm",
        "PCI compliance standards"
    ])

# ============= SLIDE 27: Future Enhancements =============
add_content_slide(prs, "Future Enhancements", [
    "Mobile App (iOS/Android)",
    "Real-time Chat Support",
    "Advanced Analytics & ML",
    "Loyalty Points System",
    "Social Media Integration",
    "Subscription Models",
    "API Marketplace",
    "Multi-language Support"
])

# ============= SLIDE 28: Performance Metrics =============
add_content_slide(prs, "Performance & Metrics", [
    "API Response Time: < 500ms average",
    "Database Query Optimization",
    "Page Load Time: < 2 seconds",
    "System Uptime: 99.9%",
    "Concurrent Users Support: 10,000+",
    "Data Security: SSL/TLS encryption",
    "Backup & Recovery: Daily backups",
    "Monitoring: Real-time alerting system"
])

# ============= SLIDE 29: Deployment Architecture =============
add_content_slide(prs, "Deployment Architecture", [
    "Frontend Hosting: CDN/Static hosting",
    "Backend Server: AWS EC2 / Linux Server",
    "Database: MySQL RDS / Self-hosted",
    "Storage: AWS S3 for media files",
    "Notifications: AWS SNS/Email Service",
    "Monitoring: CloudWatch / ELK Stack",
    "Load Balancing: Nginx / AWS ALB",
    "Caching: Redis for session management"
])

# ============= SLIDE 30: Conclusion =============
add_title_slide(prs, "THANK YOU", 
                "Svasthya Fresh - Building the Future of E-Commerce")

# Save presentation
output_path = r"c:\Users\annam\OneDrive\Documents\sf-app\Svasthya_Fresh_Presentation.pptx"
prs.save(output_path)
print(f"✓ Presentation created successfully!")
print(f"✓ Saved to: {output_path}")
print(f"✓ Total slides: {len(prs.slides)}")
